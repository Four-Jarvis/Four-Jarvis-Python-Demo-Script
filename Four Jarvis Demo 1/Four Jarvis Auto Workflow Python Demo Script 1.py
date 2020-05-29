#coding=utf-8
'''
Author:
    Jarvis
微信公众号:
    Four Jarvis
'''

import re
import sys
import os
import datetime
import time

import schedule
import docx
import openpyxl
import smtplib
import pymysql
import pandas as pd
import matplotlib.pyplot as plt
from docx.shared import RGBColor
from docx.shared import Inches
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx.util import Inches
from pptx import Presentation
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.header import Header
from selenium import webdriver
from multiprocessing import Pool
from loguru import logger

def selenium(url):
    options = webdriver.ChromeOptions()
    options.add_argument('--headless') # 无界面模式
    options.add_argument('lang=zh_CN.UTF-8') # 
    UserAgent = 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/69.0.3497.100 Safari/537.36'
    options.add_argument('User-Agent=' + UserAgent)
    prefs = { 
        'profile.default_content_setting_values': {
            'images': 2,
            'permissions.default.stylesheet':2,
            'javascript': 2
        }
    }
    options.add_experimental_option("prefs", prefs) # 优化selenium，禁止加载图片/css/javascript
    browser = webdriver.Chrome(options=options)
    browser.implicitly_wait(10) # 设置隐性等待时间为10s
#     browser.set_page_load_timeout(5)  # 设置页面加载超时
#     browser.set_script_timeout(5) # 设置页面异步js执行超时
    try:
        browser.get(url)
        pagedata = browser.page_source
        browser.quit()
        logger.info('成功访问')
        return pagedata
    except:
        browser.refresh() # 第一次若不能成功访问，则刷新页面
        time.sleep(5)
        pagedata = browser.page_source
        browser.quit()
        logger.info('刷新成功访问')
        return pagedata
    else: # 
        logger.error('无法访问')


def eastmoney2mysql(company, num):

    url = 'http://so.eastmoney.com/news/s?keyword=%s&pageindex=%s' % (company, num) # 传入所需监控的公司和爬取的页码
    
    data = selenium(url) # selenium数据采集
    p_title = '<div class="news-item"><h3><a href=".*?">(.*?)</a>' # 编写正则表达式提取数据
    p_href = '<div class="news-item"><h3><a href="(.*?)">.*?</a>'
    p_date_text = '<p class="news-desc">(.*?)</p>'

    title = re.findall(p_title,data) # findall()返回的是一个列表
    href = re.findall(p_href,data)
    date = re.findall(p_date_text,data,re.S)
    text = re.findall(p_date_text,data,re.S)
    number = len(title) # 查看标题列表的数量

    if number == 0: # 若返回的标题列表数量为0，则代表无数据，可能为东方财富网无此关键词
        logger.error(company + ' - 第' + str(num) + '页东方财富网爬取失败，无此关键词')
        return False
    else: 
        logger.info(company + ' - 第' + str(num) + '页东方财富网爬取成功')
        
        sql_1 = 'SELECT * FROM test2 WHERE company = "%s"' % company
        fm = Fmysql() # Fmysql类实例化
        data_all = fm.data_from_mysql(sql_1) # 调用Fmysql类提取数据
         # 列表推导式，存放从数据库提取有关该公司的所有标题，为录入数据库去重准备
        title_all = [data_all[j][3] for j in range(len(data_all))]
        
        try:
            for i in range(number):
                title[i] = re.sub('<.*?>', '', title[i]) # 使用正则表达式清洗数据
                date[i] = date[i].split(' ')[0]
                text[i] = re.sub('<.*?>', '', text[i], re.S)
                text[i] = re.sub('\d{4}(.*)\d{2}\s\S', '', text[i], re.S)
                text[i] = re.sub('\n', '', text[i], re.S)
                text[i] = re.sub('\s', '', text[i], re.S)

                print(str(i+1) + '.' + title[i] + ' - '+ date[i])
                print(href[i])
                print(text[i])
                # 数据去重及将数据存入数据库
                if title[i] not in title_all: # 判断数据是否在原数据库中，不在的话才进行数据存储
                    sql_2 = 'INSERT INTO test2(date, company, title, href, text) VALUES("%s", "%s", "%s", "%s", "%s")'\
                                                                    % (date[i], company, title[i], href[i], text[i])
                    fm.execute_to_mysql(sql_2)
            logger.info(company + ' - 第' + str(num) + '页东方财富网录入数据库成功')
        except:
            logger.error(company + ' - 第' + str(num) + '页东方财富网录入数据库失败')


class Fmysql:
    
    def __init__(self, host='localhost', port=3306, user='root', password='xxxxxx', database='xxxx', charset='utf8'):
        self.host = host
        self.port = port
        self.user = user
        self.password = password
        self.database = database
        self.charset = charset
        self.conn = None
        self.cur = None
        
    def connect(self): # 连接数据库
        self.conn = pymysql.connect(host=self.host, port=self.port, user=self.user, 
                                        password=self.password, database=self.database, charset=self.charset) # 跟数据库建立连接
        self.cur = self.conn.cursor()

    def close(self): # 关闭数据库
        self.cur.close()
        self.conn.close()
        
    def __execute(self, sql):
        try:
            self.connect()
            count = self.cur.execute(sql)
            self.conn.commit()
        except self.conn.Error:
            logger.error("数据库执行失败！")
            self.conn.rollback()

    def execute_to_mysql(self, sql): # 插入/修改/删除数据
        self.__execute(sql)
        self.close()

    def data_from_mysql(self, sql): # 查询数据内容
        try:
            self.__execute(sql)
            data = self.cur.fetchall()
            self.close()
        except self.conn.Error:
            logger.error("查询数据失败！")
        return data    
    
    def spec_from_mysql(self, sql): # 查询数据表头
        try:
            self.__execute(sql)
            spec = self.cur.description
            self.close()
        except self.conn.Error:
            logger.error("查询表头失败！")
        return spec


def get_Ystd():
    # 获取昨天日期的字符串格式的函数
    today = datetime.date.today() #获取今天的日期
    oneday = datetime.timedelta(days=1) #获取一天的日期格式数据
    yesterday = today - oneday
    yesterday_str = yesterday.strftime('%Y-%m-%d') #获取昨天日期的格式化字符串
    return yesterday_str #返回昨天的字符串


def get_Today():
    # 获取今天的日期的字符串格式的函数
    today = datetime.date.today()
    today_str = today.strftime('%Y-%m-%d')
    return today_str


def save2picture(date, path):
    # 从mysql读取数据并处理
    db = pymysql.connect(host='localhost', port=3306, user='root', password='xxxxxx', database='xxxx', charset='utf8')
    cursor = db.cursor()
    df = pd.read_sql('SELECT * FROM test2 WHERE date = "%s"' % date, db)
    df = df[['date', 'company']].groupby('company').count().sort_values('date', ascending = False)
    df.reset_index(inplace=True)
    mn=df['date'].min() # 当日所有公司中新闻最少的数量
    mx=df['date'].max() # 当日所有公司中新闻最多的数量

    # Draw plot
    plt.rcParams['font.sans-serif']=['SimHei'] #用来正常显示中文标签
    fig, ax = plt.subplots(figsize=(42,25), dpi= 80)
    ax.hlines(y=df.company.sort_values(ascending = False), xmin=mn, xmax=mx, color='gray', alpha=0.7, linewidth=1, linestyles='dashdot')
    ax.scatter(y=df.company.sort_values(ascending = False), x=df.date, s=75, color='firebrick', alpha=0.6)
    for date, company in enumerate(df.date.sort_values(ascending = True)): 
        ax.text(company+0.5, date, round(company, 1), horizontalalignment='center', verticalalignment='center', fontdict={'size':35})
        
    # Title, Label, Ticks and Ylim
    ax.set_title('今日新闻概览图', fontdict={'size':60})
    ax.set_xlabel('新闻数量', fontdict={'size':35})
    ax.set_yticks(df.company.sort_values(ascending = False))
    ax.set_yticklabels(df.company.str.title(), fontdict={'horizontalalignment': 'right','size':35})
    ax.set_xlim(0, mx+1)
#     plt.show()
    auto_pic = plt.savefig(path, dpi=80)
    return auto_pic


def save2excel(date, path):
    # 将数据和字段名写入excel的函数
    excel = openpyxl.Workbook() #新建一个工作薄对象
    sheet = excel.active #激活一个新的sheet
    sheet.title = date #给sheet命名
    
    #生成数据
    sql = 'SELECT * FROM test2 WHERE date = "%s"' % date
    fm = Fmysql()
    data = fm.data_from_mysql(sql) # 生成数据
    spec = fm.spec_from_mysql(sql) # 生成字段名称
    
    #导入数据
    for col in range(len(spec)): #将字段名称循环写入excel第一行，因为字段格式列表里包含列表，每个列表的第一元素才是字段名称
        #row代表行数，column代表列数，value代表单元格输入的值，行数和列数都是从1开始，这点于python不同要注意
        _ = sheet.cell(row = 1, column = col + 1, value = u'%s' % spec[col][0])
    for row in range(len(data)): #将数据循环写入excel的每个单元格中 
        for col in range(len(spec)):
            _ = sheet.cell(row = row + 2, column = col + 1, value = u'%s' % data[row][col]) #因为第一行写了字段名称，所以要从第二行开始写入
    auto_excel = excel.save(path) 
    return auto_excel #返回生成的excel


def save2docx(date, company, path):
    # 创建一个空白Word对象，并设置好字体
    word = docx.Document()
    word.styles['Normal'].font.name = u'微软雅黑'  # 可换成word里面任意字体
    word.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'微软雅黑')  # 这里也需要同时修改

    # 创建封面    
    p = word.add_picture('E:/Four Jarvis LOGO.png', width=Inches(1.25)) # 创建一个图片
    p = word.add_paragraph()  # 创建一个段落
    p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER  # 居中设置
    p.paragraph_format.space_before = Pt(150)  # 段前距为200，这个是测试出来的
    p.paragraph_format.space_after = Pt(30)  # 段后距为40，这个也是测试出来的
    run = p.add_run('Four Jarvis数据分析报告')  # 在段落里添加内容
    font = run.font  # 设置字体
    font.color.rgb = RGBColor(17, 17, 17)  # 颜色设置，这里是用RGB颜色
    font.size = Pt(36)  # 字体大小设置，和word里面的字号相对应

    p = word.add_paragraph()  # 新建一个段落
    p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(date)  # 在段落中输入当天日期
    font = run.font
    font.color.rgb = RGBColor(17, 17, 17)
    font.size = Pt(26)
    
    for x in company:
        # 添加分页符
        word.add_page_break()# 添加分页符
        
        sql = 'SELECT * FROM test2 WHERE company = "%s" AND date = "%s"' % (x, date)
        fm = Fmysql()
        data = fm.data_from_mysql(sql) # 生成数据
        num = len(data)
        
        if num != 0: 
            # 设置正文标题
            p = word.add_paragraph()
            p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER  # 段落文字居中设置
            run = p.add_run(x + '舆情报告')
            run.font.color.rgb = RGBColor(17, 17, 17)  # 字体颜色设置
            run.font.size = Pt(22)  # 字体大小设置

            # 编写正文内容之引言
            p = word.add_paragraph()  # 添加新段落
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY  # 两端对齐
            p.paragraph_format.first_line_indent = Inches(0.2)  # 控制首行缩进
            introduction = '本次舆情监控目标:%s，当天共爬取东方财富网内相关新闻%s篇，具体新闻如下：' % (x, num)
            p.add_run(introduction).bold = True # 加粗

            # 编写正文内容之具体新闻内容
            for i in range(num):
                p = word.add_paragraph()  # 添加新段落
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY  # 设置两端对齐
                p.add_run(str(i + 1) + '. ' + data[i][3])  # 提取新闻标题

            # 编写正文内容之表格添加
            tb = word.add_table(rows=num + 1, cols=2, style='Light Grid')
            tb.cell(0, 0).text = '新闻标题'
            tb.cell(0, 1).text = '新闻摘要'

            for i in range(num):
                tb.cell(i+1, 0).text = data[i][3] + data[i][4]  # 提取新闻标题 + 新闻链接
                tb.cell(i+1, 1).text = data[i][5]  # 提取新闻摘要
        else:
            # 设置正文标题
            p = word.add_paragraph()
            p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER  # 段落文字居中设置
            run = p.add_run(x + '舆情报告')
            run.font.color.rgb = RGBColor(17, 17, 17)  # 字体颜色设置
            run.font.size = Pt(22)  # 字体大小设置
            
            p = word.add_paragraph()  # 添加新段落
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY  # 两端对齐
            p.paragraph_format.first_line_indent = Inches(0.4)  # 这个控制首行缩进，
            introduction = '本次舆情监控目标:%s，当天东方财富网内无该公司相关新闻' % x
            p.add_run(introduction).bold = True # 加粗
            
    auto_docx = word.save(path)
    return auto_docx


def save2pptx(date, company, path):

    prs = Presentation('D:/python-bit/Four Jarvis/Demo1/model.pptx') # 使用预先创建好的PPT模板,也可以直接生成
    slide_1 = prs.slides[0] # 创建ppt封面 第一页
    title = slide_1.shapes[0]
    subtitle = slide_1.shapes[1]
    title.text = "Four Jarvis 数据分析报告"
    subtitle.text = date
    
    img_path = 'D:/python-bit/Four Jarvis/Demo1/Four Jarvis Auto Demo1 %s.jpg' % date # 图片储存路径
    slide_2 = prs.slides[1] # 创建ppt 第二页 插入图片概览
    pic = slide_2.shapes.add_picture(img_path, left=Inches(0), top=Inches(0), height=Inches(8)) 
    
    for x in company:        
        #templateStyleNum = len(prs.slide_layouts) # 获取模板个数
        oneSlide = prs.slides.add_slide(prs.slide_layouts[0]) # 按照第1个模板创建 一张幻灯片 
        sql = 'SELECT * FROM test2 WHERE company = "%s" AND date = "%s"' % (x, date)
        fm = Fmysql()
        data = fm.data_from_mysql(sql) # 生成数据
        num = len(data)
        
        if num != 0:        
            body_shapes = oneSlide.shapes.placeholders # 获取模板可填充的所有位置
            table_placeholder = oneSlide.shapes[2]
            
            for index, body_shape in enumerate(body_shapes):
                if index == 0:
                    body_shape.text = '%s舆情报告' % x
                elif index == 1:
                    body_shape.text = '本次舆情监控目标:%s，当天共爬取东方财富网内相关新闻%s篇，具体新闻如下：'% (x, num)
                elif index == 2:
                    rows, cols = num + 1, 2
                    tb = table_placeholder.insert_table(rows, cols).table  # 添加表格，并取表格类

                    tb.columns[0].width = Inches(4.5)  # 第一列宽度
                    tb.columns[1].width = Inches(8.8)  # 第二列宽度
                    tb.cell(0, 0).text = '新闻标题'
                    tb.cell(0, 1).text = '新闻摘要'

                    for i in range(num):
                        tb.cell(i+1, 0).text = data[i][3] + data[i][4]  # 提取新闻标题 + 新闻链接               
                        tb.cell(i+1, 1).text = data[i][5]  # 提取新闻摘要                        
        else:
            body_shapes = oneSlide.shapes.placeholders
            for index, body_shape in enumerate(body_shapes):
                if index == 0:
                    body_shape.text = '%s舆情报告' % x
                elif index == 1:
                    body_shape.text = '本次舆情监控目标:%s，当天东方财富网内无该公司相关新闻!' % x
    auto_pptx = prs.save(path)
    return auto_pptx


def create_email(email_from, email_to, email_subject, email_text, attach_path, attach_name):
    # 输入发件人昵称、收件人昵称、主题，正文，附件地址,附件名称生成一封邮件
    message = MIMEMultipart() # 生成一个空的带附件的邮件实例
    message.attach(MIMEText(email_text, 'plain', 'utf-8')) # 将正文以text的形式插入邮件中
    message['From'] = Header(email_from, 'utf-8') # 生成发件人名称（这个跟发送的邮件没有关系）
    message['To'] = Header(email_to, 'utf-8') # 生成收件人名称（这个跟接收的邮件也没有关系）
    message['Subject'] = Header(email_subject, 'utf-8') # 生成邮件主题                 
    for name in attach_name:
        if os.path.isfile(attach_path + '/' + name):
            # 构造附件
            att = MIMEText(open(attach_path + '/' + name, 'rb').read(), 'base64', 'utf-8')
            att["Content-Type"] = 'application/octet-stream'
            att.add_header("Content-Disposition", "attachment", filename=("gbk", "", name))
            message.attach(att)
    return message # 返回邮件


def send_email(sender, password, receiver, message):
    # 输入发件人邮箱、密码、收件人邮箱、邮件内容发送邮件的函数
    try:
        server = smtplib.SMTP_SSL("smtp.qq.com", 465)  # 发件人邮箱中的SMTP服务器
        server.ehlo()
        server.login(sender, password)  # 登录qq邮箱
        #发送邮件
        server.sendmail(sender, receiver, message.as_string())  # 括号中对应的是发件人邮箱账号、收件人邮箱账号（是一个列表）、邮件内容
        print("邮件发送成功")
        server.quit()  # 关闭连接
    except Exception:
        print(traceback.print_exc())
        print("邮件发送失败")


def main_by_schedule(main):
    schedule.every().day.at("23:43").do(main)        #每天23:43执行一次  
#     schedule.every(5).minutes.do(main)              #每5分钟执行一次
#     schedule.every().hour.do(main)                   #每小时执行一次  
#     schedule.every().monday.do(main)                 #每周一执行一次
#     schedule.every().monday.at("22:45").do(main)  #每周一22:45执行一次
   
    while True:
        schedule.run_pending()
        time.sleep(7) 
        # 这里时间设置是为schedule函数停顿一下，可改为1s，while True会一直循环执行


def main():
    
    # 监控目标公司列表，仅作演示
    companies = ['高伟达', '飞天诚信', '朗科科技', '永太科技', '海螺水泥', '瑞幸咖啡', '腾讯控股', '腾讯', '分众传媒',\
            '阿里巴巴', '中国联通', 'TCL科技', '中公教育', '华北制药', '亿帆医药', '圣达生物', '振华股份', '三一重工',\
            '海天味业', '贵州茅台', '美的集团', '中国平安', '华润三九', '索菲亚', '格力电器', '白云机场', '海康威视']
    page = 2 # 每次爬取页面数量
    path = 'D:/python-bit/Four Jarvis/Demo1' # 文件存储路径
    today = get_Today() # 或者自定义日期'xxxx-xx-xx'
    logger.add(f'Four Jarvis Auto Demo1 {today}.log') # 日志记录
    
    print('爬虫程序开始执行：')
    start = time.time()    
    #单机串行实现
#     for num in range(page):
#         for company in companies:
#             eastmoney2mysql(company, num + 1)
            
#     多进程实现数据采集和存储
    pool = Pool(3) # 创建进程池
    for num in range(page):
        for company in companies:
            pool.apply_async(eastmoney2mysql(company, num + 1))
    pool.close()
    pool.join()

    end = time.time()
    print('爬虫程序执行结束！总用时：' + str(end-start))
    
    #AUTO_PICTURE  
    my_file_name1 = 'Four Jarvis Auto Demo1 %s.jpg' % today
    file_path1 = path + '/' + my_file_name1
    save2picture(today, file_path1) # 生成picture    
    
    #AUTO_EXCEL  
    my_file_name2 = 'Four Jarvis Auto Demo1 %s.xlsx' % today
    file_path2 = path + '/' + my_file_name2
    save2excel(today, file_path2) # 生成excel
    
    #AUTO_WORD
    my_file_name3 = 'Four Jarvis Auto Demo1 %s.docx' % today
    file_path3 = path + '/' + my_file_name3
    save2docx(today, companies, file_path3) # 生成docx
    
    #AUTO_PPT
    my_file_name4 = 'Four Jarvis Auto Demo1 %s.pptx' % today
    file_path4 = path + '/' + my_file_name4
    save2pptx(today, companies, file_path4) # 生成pptx

    #AUTO_MAIL
    my_email_from = 'Four Jarvis'
    my_email_to = 'User'
    my_email_subject = 'Four Jarvis Auto Demo1 ' + today # 邮件标题
    my_email_text = "Dear all,\n\t\t附件为今天的数据分析报告，请查收！\n\nFour Jarvis " # 邮件正文
    
    my_attach_name = [my_file_name1, my_file_name2, my_file_name3, my_file_name4] # 附件名称
    my_message = create_email(my_email_from, my_email_to, my_email_subject, my_email_text, path, my_attach_name)  # 生成邮件
    my_sender = 'xxxxxxxxx@qq.com'
    my_password = 'xxxxxxxxxxxxxxxxx' # QQ邮箱的SMTP授权码
    my_receiver = ['xxxxxxxxx@qq.com', 'xxxxxxxxx@qq.com', 'xxxxxxxxx@qq.com'] # 接收人邮箱列表
    send_email(my_sender, my_password, my_receiver, message) # 发送邮件

if __name__ == '__main__':
    main_by_schedule(main) # 定时执行

